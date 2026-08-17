"""Trích xuất văn bản và metadata trang từ tài liệu."""

import logging
import os


def _extract_pdf_pages(filepath):
    """Trích xuất từng trang PDF để giữ số trang cho phần trích dẫn."""
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer

    pages = []
    for page_number, page_layout in enumerate(extract_pages(filepath), start=1):
        text = "".join(
            element.get_text()
            for element in page_layout
            if isinstance(element, LTTextContainer)
        ).strip()
        if text:
            pages.append({"page": page_number, "text": text})
    return pages


def extract_text_by_page(filepath):
    """Trả về danh sách ``{"page": số_trang, "text": nội_dung}``.

    PDF được tách theo trang. Các định dạng không có khái niệm trang được trả về
    như một phần tử duy nhất với ``page=None``.
    """
    if os.path.splitext(filepath)[1].lower() == ".pdf":
        return _extract_pdf_pages(filepath)

    text = extract_text(filepath)
    return [{"page": None, "text": text}] if text else []


def extract_text(filepath):
    """Trích xuất văn bản từ PDF, Word, Excel, PowerPoint, TXT hoặc Markdown."""
    file_ext = os.path.splitext(filepath)[1].lower()

    if file_ext == ".pdf":
        return "\n\n".join(page["text"] for page in _extract_pdf_pages(filepath))

    if file_ext in (".txt", ".md"):
        with open(filepath, "r", encoding="utf-8") as file:
            return file.read()

    if file_ext == ".docx":
        try:
            from docx import Document

            document = Document(filepath)
            return "\n".join(paragraph.text for paragraph in document.paragraphs)
        except ImportError:
            logging.error("Cần cài python-docx để đọc tài liệu Word")
            return ""

    if file_ext in (".xlsx", ".xls"):
        try:
            import pandas as pd

            text_parts = []
            workbook = pd.ExcelFile(filepath)
            for sheet_name in workbook.sheet_names:
                dataframe = workbook.parse(sheet_name)
                text_parts.append(f"Trang tính: {sheet_name}\n")
                text_parts.append(dataframe.to_string(index=False))
            return "\n\n".join(text_parts)
        except ImportError:
            logging.error("Cần cài pandas để đọc bảng tính Excel")
            return ""

    if file_ext == ".pptx":
        try:
            from pptx import Presentation

            presentation = Presentation(filepath)
            text_parts = []
            for slide_number, slide in enumerate(presentation.slides, start=1):
                text_parts.append(f"Trang chiếu {slide_number}")
                text_parts.extend(
                    shape.text for shape in slide.shapes if hasattr(shape, "text")
                )
            return "\n".join(text_parts)
        except ImportError:
            logging.error("Cần cài python-pptx để đọc tệp PowerPoint")
            return ""

    logging.warning("Định dạng tệp chưa được hỗ trợ: %s", file_ext)
    return ""
